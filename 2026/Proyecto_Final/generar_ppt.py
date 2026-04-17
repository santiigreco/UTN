"""
Genera la presentación PPT del Proyecto Final G8 - Planta Sin TACC
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import os, io

# ─── COLORES ───
FOREST  = RGBColor(0x2d, 0x5a, 0x3d)
NAVY    = RGBColor(0x2c, 0x3e, 0x6b)
TERRA   = RGBColor(0xb8, 0x5c, 0x3c)
WHEAT   = RGBColor(0xc8, 0xa9, 0x6e)
PLUM    = RGBColor(0x6b, 0x3a, 0x6b)
WHITE   = RGBColor(0xff, 0xff, 0xff)
BLACK   = RGBColor(0x1a, 0x1a, 0x1a)
GRAY    = RGBColor(0x4a, 0x45, 0x40)
LGRAY   = RGBColor(0x8a, 0x84, 0x7a)
BG      = RGBColor(0xfa, 0xf9, 0xf6)
SURF    = RGBColor(0xf3, 0xf1, 0xec)

# matplotlib colors
C_FOREST = '#2d5a3d'
C_NAVY   = '#2c3e6b'
C_TERRA  = '#b85c3c'
C_WHEAT  = '#c8a96e'
C_PLUM   = '#6b3a6b'
C_SAGE   = '#7a9a6a'
C_SLATE  = '#5a6a7a'

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── HELPERS ───
def solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, fill_color=None, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shp.shadow.inherit = False
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, left, top, w, h, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_frame(slide, left, top, w, h, items, size=13, color=GRAY, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.space_before = Pt(0)
        # Split bold prefix if present
        if '→' in item:
            parts = item.split('→', 1)
            run = p.add_run()
            run.text = parts[0].strip()
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.color.rgb = BLACK
            run.font.name = 'Calibri'
            run2 = p.add_run()
            run2.text = ' — ' + parts[1].strip()
            run2.font.size = Pt(size)
            run2.font.color.rgb = color
            run2.font.name = 'Calibri'
        else:
            run = p.add_run()
            run.text = '  ' + item
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = 'Calibri'
        p.bullet = True
    return txBox

def add_kpi(slide, left, top, w, value, label, val_color=FOREST, sub=None):
    h = Inches(1.15)
    card = add_shape(slide, left, top, w, h, fill_color=WHITE, line_color=SURF)
    add_text(slide, left + Inches(0.15), top + Inches(0.12), w - Inches(0.3), Inches(0.3),
             label, size=9, bold=True, color=LGRAY, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + Inches(0.38), w - Inches(0.2), Inches(0.5),
             value, size=26, bold=True, color=val_color, align=PP_ALIGN.CENTER, font_name='Consolas')
    if sub:
        add_text(slide, left + Inches(0.1), top + Inches(0.85), w - Inches(0.2), Inches(0.25),
                 sub, size=9, color=LGRAY, align=PP_ALIGN.CENTER)
    return card

def section_header(slide, num, title, lead=""):
    add_shape(slide, Inches(0.6), Inches(0.4), Inches(0.45), Inches(0.45), fill_color=FOREST)
    add_text(slide, Inches(0.6), Inches(0.42), Inches(0.45), Inches(0.45),
             str(num), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.2), Inches(0.35), Inches(8), Inches(0.5),
             title, size=26, bold=True, color=BLACK)
    if lead:
        add_text(slide, Inches(1.2), Inches(0.85), Inches(10), Inches(0.45),
                 lead, size=13, color=GRAY)

def make_chart_img(fig, dpi=180):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight', facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf

def add_table(slide, left, top, w, rows_data, col_widths, header_color=FOREST):
    rows = len(rows_data)
    cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, w, Inches(0.38 * rows))
    tbl = tbl_shape.table
    # Set column widths
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw
    # Style
    for r_idx, row_data in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = 'Calibri'
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.size = Pt(9)
                else:
                    paragraph.font.color.rgb = GRAY
                # Right-align numeric columns
                if c_idx > 0 and r_idx > 0:
                    paragraph.alignment = PP_ALIGN.RIGHT
                if c_idx > 0 and r_idx == 0:
                    paragraph.alignment = PP_ALIGN.RIGHT
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r_idx == len(rows_data) - 1 and row_data[0].startswith('Total') or row_data[0].startswith('TOTAL'):
                cell.fill.solid()
                cell.fill.fore_color.rgb = SURF
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = BLACK
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return tbl_shape


# ══════════════════════════════════════════════════════════════
# SLIDE 0 — PORTADA
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
solid_bg(sl, FOREST)

# Overlay darker strip at bottom
add_shape(sl, Inches(0), Inches(4.5), Inches(13.333), Inches(3), fill_color=RGBColor(0x1a, 0x3a, 0x28))

add_text(sl, Inches(0.8), Inches(0.6), Inches(4), Inches(0.35),
         'UTN FRBA — GRUPO 8', size=12, bold=True, color=WHEAT, font_name='Consolas')

add_text(sl, Inches(0.8), Inches(1.4), Inches(10), Inches(1.5),
         'Planta Industrial de Alimentos\nLibres de Gluten', size=42, bold=True, color=WHITE)

add_text(sl, Inches(0.8), Inches(3.3), Inches(8), Inches(0.6),
         'Estudio de prefactibilidad para la primera planta 100% aislada de Argentina\ndedicada a la produccion de alimentos sin TACC a escala industrial.',
         size=16, color=RGBColor(0xcc, 0xcc, 0xcc))

# KPIs en la franja oscura
kpi_data = [('38,85%', 'TIR'), ('$1.964 MM', 'INVERSION'), ('3,5 años', 'PAYBACK'), ('7 SKUs', 'PRODUCTOS')]
for i, (v, l) in enumerate(kpi_data):
    x = Inches(0.8 + i * 2.8)
    add_text(sl, x, Inches(5.0), Inches(2.2), Inches(0.6), v, size=32, bold=True, color=WHEAT, font_name='Consolas')
    add_text(sl, x, Inches(5.65), Inches(2.2), Inches(0.3), l, size=10, bold=True, color=LGRAY, font_name='Consolas')

add_text(sl, Inches(0.8), Inches(6.7), Inches(10), Inches(0.3),
         'Albarracin  ·  Carrizo  ·  Greco  ·  Rodriguez  ·  Schaltegger', size=11, color=LGRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — MERCADO OBJETIVO
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG)
section_header(sl, 1, 'Mercado Objetivo',
    'Un mercado de +1,3M de consumidores que paga hasta 3x mas por productos de baja calidad.')

# KPIs
kpis = [('46M', 'Poblacion AR', NAVY), ('460.000', 'Celiacos (1%)', TERRA),
        ('1,33M', 'Mercado ampliado', FOREST), ('24 kg/año', 'Consumo per capita', WHEAT)]
for i, (v, l, c) in enumerate(kpis):
    add_kpi(sl, Inches(0.6 + i * 3.05), Inches(1.5), Inches(2.8), v, l, val_color=c)

# Chart — Donut mercado
fig, ax = plt.subplots(figsize=(4.5, 3.5))
sizes = [460, 230, 644]
labels = ['Celiacos\ndiagnosticados\n(460K)', 'Convivientes\n(+230K)', 'Consumo\nSaludable\n(+644K)']
colors = [C_TERRA, C_WHEAT, C_FOREST]
wedges, texts, autotexts = ax.pie(sizes, labels=labels, colors=colors, autopct='%1.0f%%',
    startangle=90, pctdistance=0.75, wedgeprops=dict(width=0.4, edgecolor='white', linewidth=2))
for t in texts: t.set_fontsize(9)
for t in autotexts: t.set_fontsize(9); t.set_color('white'); t.set_weight('bold')
ax.set_title('Composicion del mercado target', fontsize=11, fontweight='bold', color='#4a4540', pad=12)
chart_buf = make_chart_img(fig)
sl.shapes.add_picture(chart_buf, Inches(0.6), Inches(2.9), Inches(5), Inches(3.8))

# Bullets oportunidad
items = [
    'Celiaquía→afecta al 1% de la poblacion (fuente ANMAT). Diagnostico crece cada año.',
    'Convivientes→adaptan su dieta al hogar libre de gluten, multiplican base x1,5.',
    'Consumo consciente→veganos, plant-based, clean label. Factor adicional x1,4.',
    'Market share objetivo→7% del mercado ampliado (~93.380 consumidores). Conservador vs SMAMS (~10%).',
    'Tendencia global→el mercado gluten-free crece sostenidamente (ref: SA GF Market Report 2030).',
]
add_bullet_frame(sl, Inches(6), Inches(2.9), Inches(6.8), Inches(4.2), items, size=12)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — PRODUCTO Y DIFERENCIACION
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG)
section_header(sl, 2, 'Producto y Diferenciacion',
    '7 SKUs industriales, veganos, sin alergenos. Precios entre $7.000 y $13.700/kg vs. referencia $17.500/kg.')

# Tabla de productos
prod_data = [
    ['Producto', 'Linea', 'Descripcion', 'Precio/kg'],
    ['Chocolate con chips picante', 'A', 'Galletita dulce con chocolate y especias', '$13.700'],
    ['Galletitas simil Chocolinas', 'A', 'Version sin TACC, endulzada naturalmente', '$10.500'],
    ['Pepas chocolate y batata', 'A', 'Base suave con relleno de dulce de batata', '$11.000'],
    ['Bizcochos salados', 'A', 'Snack salado de harinas alternativas', '$7.000'],
    ['Polvorones', 'A', 'Desmenuzables: vainilla, chocolate, limon', '$10.500'],
    ['Pan de molde', 'B', 'Pan lactal rebanado, alta rotacion', '$12.700'],
    ['Budines', 'B', 'Saborizados (vainilla, chocolate, frutas)', '$12.000'],
]
add_table(sl, Inches(0.6), Inches(1.5), Inches(7.5), prod_data,
          [Inches(2.4), Inches(0.6), Inches(3.2), Inches(1.3)])

# Chart — Precios vs referencia
fig, ax = plt.subplots(figsize=(4.8, 3.8))
prods = ['Choc. chips', 'Pan molde', 'Budines', 'Pepas', 'Chocolinas', 'Polvorones', 'Bizcochos']
precios = [13700, 12700, 12000, 11000, 10500, 10500, 7000]
bars = ax.barh(prods, precios, color=C_FOREST, height=0.55, zorder=3)
ax.axvline(x=17500, color=C_WHEAT, linestyle='--', linewidth=2, label='Ref. mercado $17.500/kg')
ax.set_xlim(0, 20000)
ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'${int(x/1000)}K'))
ax.legend(fontsize=9, loc='lower right')
ax.set_title('Nuestro precio vs. referencia de mercado', fontsize=11, fontweight='bold', color='#4a4540')
ax.invert_yaxis()
ax.grid(axis='x', alpha=0.2)
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
for bar, p in zip(bars, precios):
    ax.text(bar.get_width() + 200, bar.get_y() + bar.get_height()/2, f'${p:,}'.replace(',','.'), va='center', fontsize=9, color='#4a4540')
chart_buf = make_chart_img(fig)
sl.shapes.add_picture(chart_buf, Inches(8.3), Inches(1.5), Inches(4.6), Inches(3.6))

# Diferenciacion bullets
items = [
    'Planta 100% aislada→cero contaminacion cruzada por diseño, no por protocolo.',
    'Vegano + sin alergenos→amplifica base de consumidores.',
    'Endulzado natural→alineado con tendencia clean label.',
    'Certificaciones→FSSC 22000 + HACCP + ANMAT + Logo SIN TACC + Kosher.',
]
add_bullet_frame(sl, Inches(0.6), Inches(5.3), Inches(12), Inches(2), items, size=12)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — PLAN DE VENTAS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG)
section_header(sl, 3, 'Plan de Ventas',
    '383,5 toneladas/año en 7 SKUs a traves de 4 canales. Facturacion anual $4.382 MM.')

kpis = [('383,5 tn', 'Produccion anual', FOREST), ('$4.382M', 'Facturacion anual', NAVY), ('$43.819M', 'Acum. 10 años', PLUM)]
for i, (v, l, c) in enumerate(kpis):
    add_kpi(sl, Inches(0.6 + i * 4.1), Inches(1.5), Inches(3.8), v, l, val_color=c)

# Chart — Revenue por producto
fig, ax = plt.subplots(figsize=(5.5, 3.5))
prods = ['Pan molde', 'Pepas', 'Budines', 'Choc. chips', 'Polvorones', 'Chocolinas', 'Bizcochos']
revs = [952.5, 858, 849.6, 657.6, 554.4, 291.4, 218.4]
palette = [C_FOREST, C_NAVY, C_TERRA, C_WHEAT, C_PLUM, C_SAGE, C_SLATE]
bars = ax.barh(prods, revs, color=palette, height=0.6, zorder=3)
ax.set_xlim(0, 1100)
ax.set_title('Revenue anual por producto (ARS MM)', fontsize=11, fontweight='bold', color='#4a4540')
ax.invert_yaxis()
ax.grid(axis='x', alpha=0.2)
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
for bar, r in zip(bars, revs):
    ax.text(bar.get_width() + 15, bar.get_y() + bar.get_height()/2, f'${r:,.0f}M'.replace(',','.'), va='center', fontsize=9, color='#4a4540')
chart_buf = make_chart_img(fig)
sl.shapes.add_picture(chart_buf, Inches(0.6), Inches(2.9), Inches(6), Inches(3.8))

# Canales
items = [
    'Distribuidores mayoristas→canal principal (B2B), cobertura nacional.',
    'Supermercados e hipermercados→gondola de productos saludables/dieteticos.',
    'Dieteticas→canal especializado, alto margen, afinidad celiaco.',
    'Farmacias→canal complementario para productos de salud.',
]
add_bullet_frame(sl, Inches(7), Inches(3.0), Inches(5.8), Inches(3.5), items, size=12)

# Tabla ventas
ventas_data = [
    ['Producto', 'Tn/año', 'Precio/Tn', 'Revenue (MM)', '%'],
    ['Pan de molde', '75,0', '$12.700', '$952,5', '21,7%'],
    ['Pepas choc. y batata', '78,0', '$11.000', '$858,0', '19,6%'],
    ['Budines', '70,8', '$12.000', '$849,6', '19,4%'],
    ['Choc. chips picante', '48,0', '$13.700', '$657,6', '15,0%'],
    ['Polvorones', '52,8', '$10.500', '$554,4', '12,7%'],
    ['Chocolinas', '27,75', '$10.500', '$291,4', '6,7%'],
    ['Bizcochos salados', '31,2', '$7.000', '$218,4', '5,0%'],
    ['TOTAL', '383,55', '—', '$4.381,9', '100%'],
]
add_table(sl, Inches(7), Inches(4.7), Inches(5.8), ventas_data,
          [Inches(1.7), Inches(0.8), Inches(0.9), Inches(1.2), Inches(0.7)])


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — PROCESO PRODUCTIVO
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG)
section_header(sl, 4, 'Proceso Productivo',
    '10 etapas en dos lineas dedicadas. Horno tunel continuo + hornos rotativos. Detector de metales integrado.')

# Flujo — 10 pasos en 2 filas de 5
steps = ['1. Recepcion\nde MP', '2. Almacena-\nmiento', '3. Pesado y\ndosificado', '4. Mezclado/\namasado', '5. Moldeado/\nformado',
         '6. Horneado', '7. Enfriado', '8. Envasado +\ndet. metales', '9. Etiquetado', '10. Almacen\nPT']
for i, step in enumerate(steps):
    row = i // 5
    col = i % 5
    x = Inches(0.6 + col * 2.4)
    y = Inches(1.5 + row * 1.15)
    card = add_shape(sl, x, y, Inches(2.15), Inches(0.95), fill_color=WHITE, line_color=SURF)
    add_text(sl, x, y + Inches(0.15), Inches(2.15), Inches(0.7),
             step, size=11, bold=True, color=FOREST, align=PP_ALIGN.CENTER)

# Linea A
add_shape(sl, Inches(0.6), Inches(3.9), Inches(5.8), Inches(3.2), fill_color=WHITE, line_color=SURF)
add_text(sl, Inches(0.8), Inches(4.0), Inches(5), Inches(0.3),
         'LINEA A — GALLETITAS Y SNACKS', size=11, bold=True, color=NAVY)
add_text(sl, Inches(0.8), Inches(4.3), Inches(5), Inches(0.3),
         'Capacidad: 100 kg/h  ·  Utilizacion: 61,9%', size=10, color=LGRAY, font_name='Consolas')
items_a = ['Amasadora industrial (USD 13.500)', 'Rotativa con dulcera (USD 15.000)',
           '4 moldes formadores (USD 5.000 c/u)', 'Horno tunel continuo (USD 100.000)',
           'Cinta de enfriamiento por aire forzado', 'Envasadora vertical + detector de metales']
add_bullet_frame(sl, Inches(0.8), Inches(4.7), Inches(5.4), Inches(2.3), items_a, size=11, color=GRAY)

# Linea B
add_shape(sl, Inches(6.8), Inches(3.9), Inches(5.8), Inches(3.2), fill_color=WHITE, line_color=SURF)
add_text(sl, Inches(7.0), Inches(4.0), Inches(5), Inches(0.3),
         'LINEA B — PANIFICADOS Y BUDINES', size=11, bold=True, color=TERRA)
add_text(sl, Inches(7.0), Inches(4.3), Inches(5), Inches(0.3),
         'Capacidad: 50 kg/h  ·  Utilizacion: 75,9%', size=10, color=LGRAY, font_name='Consolas')
items_b = ['Batidora industrial con accesorio batidor', 'Dosificadora volumetrica de pistones',
           '2 hornos rotativos con carros ($30 MM)', 'Fermentadora de pan + cortadora (20 reb.)',
           'Detector de metales obligatorio', 'Envasadora Flow Pack (USD 50.000)']
add_bullet_frame(sl, Inches(7.0), Inches(4.7), Inches(5.4), Inches(2.3), items_b, size=11, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — LOCALIZACION (MACRO + MICRO)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG)
section_header(sl, 5, 'Localizacion e Infraestructura',
    'Macrolocalizacion (3 regiones) y microlocalizacion (3 parques industriales, matriz ponderada).')

# KPIs
add_kpi(sl, Inches(0.6), Inches(1.5), Inches(2.8), '1.330 m²', 'Superficie planta', val_color=PLUM)
add_kpi(sl, Inches(3.6), Inches(1.5), Inches(3.2), 'P.I. Garin', 'Ubicacion elegida', val_color=NAVY, sub='GBA Norte · Panamericana')
add_kpi(sl, Inches(7.0), Inches(1.5), Inches(2.8), '8,45/10', 'Puntaje micro', val_color=FOREST, sub='1ro de 3 candidatos')

# Macro tabla
macro_data = [
    ['Region', 'Ventaja', 'Desventaja'],
    ['A — Centro (Cordoba)', 'Cercanía a fuentes de harina sin TACC', 'Lejos del mercado de consumo'],
    ['B — Litoral (Santa Fe)', 'Punto logistico para recepcion MP', 'Mercado consumidor diluido'],
    ['C — AMBA / GBA  ✓', '~40% poblacion, centros de distribucion', '—'],
]
add_table(sl, Inches(0.6), Inches(2.9), Inches(5.8), macro_data,
          [Inches(1.8), Inches(2.4), Inches(1.6)], header_color=NAVY)

add_text(sl, Inches(0.6), Inches(4.6), Inches(5.8), Inches(0.6),
         'Decision: Region C (AMBA/GBA). Los centros de distribucion de las principales cadenas estan en GBA. Es mas barato traer MP a granel que despachar PT.',
         size=10, color=GRAY)

# Micro tabla
micro_data = [
    ['Factor', 'Peso', 'Ezeiza', 'Garin', 'Pilar'],
    ['Proximidad distrib./clientes', '0,30', '7', '8', '6'],
    ['Infraestructura y servicios', '0,25', '8', '9', '9'],
    ['Marco legal y zonificacion', '0,20', '9', '9', '9'],
    ['Costo superficie y expansion', '0,15', '9', '8', '9'],
    ['Mano de obra e infra. vial', '0,10', '7', '8', '7'],
    ['TOTAL PONDERADO', '1,00', '7,95', '8,45', '7,90'],
]
add_table(sl, Inches(6.8), Inches(2.9), Inches(6), micro_data,
          [Inches(2.2), Inches(0.7), Inches(0.8), Inches(0.8), Inches(0.8)], header_color=FOREST)

# Score bars as chart
fig, ax = plt.subplots(figsize=(5, 1.5))
parks = ['P.I. Pilar', 'P.I. Ezeiza', 'P.I. Garin']
scores = [7.90, 7.95, 8.45]
colors_bars = [C_TERRA, C_WHEAT, C_FOREST]
bars = ax.barh(parks, scores, color=colors_bars, height=0.5, zorder=3)
ax.set_xlim(7, 9)
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
for bar, s in zip(bars, scores):
    ax.text(bar.get_width() + 0.03, bar.get_y() + bar.get_height()/2, f'{s:.2f}', va='center', fontsize=10, fontweight='bold', color='#4a4540')
ax.grid(axis='x', alpha=0.2)
chart_buf = make_chart_img(fig)
sl.shapes.add_picture(chart_buf, Inches(6.8), Inches(5.5), Inches(5.5), Inches(1.6))


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — COSTOS PRINCIPALES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG)
section_header(sl, 6, 'Costos Principales',
    'Costo operativo $3.060 MM/año. Margen bruto 30,2%. Margen neto 18,1%.')

kpis = [('$3.060M', 'Costo total anual', TERRA), ('$1.322M', 'Resultado bruto', FOREST, '30,2%'),
        ('$793M', 'Resultado neto', NAVY, '18,1%')]
for i, kd in enumerate(kpis):
    v, l, c = kd[0], kd[1], kd[2]
    sub = kd[3] if len(kd) > 3 else None
    add_kpi(sl, Inches(0.6 + i * 4.1), Inches(1.5), Inches(3.8), v, l, val_color=c, sub=sub)

# Chart — Costos donut
fig, ax = plt.subplots(figsize=(4.5, 3.5))
csizes = [1197.5, 944.8, 918.1]
clabels = ['Mano de obra\n$1.197M (39%)', 'Materia prima\n$945M (31%)', 'Otros costos\n$918M (30%)']
ccolors = [C_NAVY, C_FOREST, C_WHEAT]
wedges, texts, autotexts = ax.pie(csizes, labels=clabels, colors=ccolors, autopct='%1.0f%%',
    startangle=90, pctdistance=0.75, wedgeprops=dict(width=0.4, edgecolor='white', linewidth=2))
for t in texts: t.set_fontsize(9)
for t in autotexts: t.set_fontsize(9); t.set_color('white'); t.set_weight('bold')
ax.set_title('Composicion de costos operativos', fontsize=11, fontweight='bold', color='#4a4540', pad=12)
chart_buf = make_chart_img(fig)
sl.shapes.add_picture(chart_buf, Inches(0.6), Inches(2.9), Inches(5), Inches(3.8))

# Tabla personal
personal_data = [
    ['Puesto', 'Cant.', 'Sueldo/mes', 'Costo anual (MM)'],
    ['Operarios de produccion', '32', '$1.072.000', '$510,4'],
    ['Mantenimiento y logistica', '6', '$1.395.200', '$124,6'],
    ['Staff de produccion', '6', '$2.500.000', '$223,2'],
    ['Administrativos', '4', '$2.200.000', '$130,9'],
    ['Gerentes (planta + admin)', '2', '$7.000.000', '$208,3'],
    ['TOTAL', '50', '—', '$1.197,5'],
]
add_table(sl, Inches(6.2), Inches(2.9), Inches(6.5), personal_data,
          [Inches(2.3), Inches(0.7), Inches(1.4), Inches(1.5)])

# Cascada
add_shape(sl, Inches(6.2), Inches(5.9), Inches(6.5), Inches(1.1), fill_color=RGBColor(0xe6, 0xf0, 0xea), line_color=FOREST)
add_text(sl, Inches(6.4), Inches(6.0), Inches(6.1), Inches(0.9),
         'Facturacion $4.382M − Costos $3.060M = Res. bruto $1.322M (30,2%)\n− Hon. directorio 5% ($66M) − IIGG 35% ($463M) = Resultado neto $793M (18,1%)',
         size=11, bold=False, color=FOREST)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — ORGANIGRAMA
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG)
section_header(sl, 7, 'Organigrama',
    'Estructura de 50 personas con doble control de calidad y produccion en 3 turnos.')

# Nivel 1
add_shape(sl, Inches(5.5), Inches(1.6), Inches(2.3), Inches(0.7), fill_color=FOREST)
add_text(sl, Inches(5.5), Inches(1.7), Inches(2.3), Inches(0.5), 'Presidente', size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_shape(sl, Inches(5.5), Inches(2.6), Inches(2.3), Inches(0.9), fill_color=FOREST)
add_text(sl, Inches(5.5), Inches(2.65), Inches(2.3), Inches(0.3), 'Gerente General', size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, Inches(5.5), Inches(2.95), Inches(2.3), Inches(0.3), '+ Asesoria legal/contable', size=9, color=RGBColor(0xcc,0xcc,0xcc), align=PP_ALIGN.CENTER)

# Nivel 2
areas = [('Ventas', 'Vendedores'), ('Gte.\nProduccion', '~40 pers.'), ('Aseg.\nCalidad', '3 analistas'),
         ('Calidad\ne I+D', '3 inspectores'), ('Planeamiento', '3 operadores'), ('Admin y\nFinanzas', '4 areas')]
for i, (title, detail) in enumerate(areas):
    x = Inches(0.3 + i * 2.1)
    card = add_shape(sl, x, Inches(4.0), Inches(1.9), Inches(1.0), fill_color=WHITE, line_color=WHEAT)
    add_text(sl, x, Inches(4.1), Inches(1.9), Inches(0.5), title, size=11, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(sl, x, Inches(4.6), Inches(1.9), Inches(0.3), detail, size=9, color=LGRAY, align=PP_ALIGN.CENTER)

# Detalles
details = [
    ('Produccion', ['Jefe Mant.: 1 sup. + 3 tecnicos', '3 turnos: M(30), T(20), N(12)', 'SySH: 4 operadores']),
    ('Calidad (doble control)', ['Aseguramiento: 1 calidad + 1 proveedores + 2 QA', 'Inspeccion: 3 inspectores in-line', 'Clave para FSSC 22000 y HACCP']),
    ('Administracion', ['Cuentas a Pagar', 'Tesoreria', 'Facturacion y Cobranzas', 'RRHH (50 personas)']),
]
for i, (title, items) in enumerate(details):
    x = Inches(0.6 + i * 4.2)
    add_shape(sl, x, Inches(5.4), Inches(3.9), Inches(1.8), fill_color=WHITE, line_color=SURF)
    add_text(sl, x + Inches(0.15), Inches(5.5), Inches(3.6), Inches(0.3),
             title, size=10, bold=True, color=FOREST)
    add_bullet_frame(sl, x + Inches(0.15), Inches(5.8), Inches(3.6), Inches(1.3),
                    items, size=10, color=GRAY, spacing=Pt(3))


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — INVERSIONES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG)
section_header(sl, 8, 'Inversiones',
    '$1.964 MM (~USD 1,4M) con 71% destinado a la nave industrial.')

# KPI grande
add_shape(sl, Inches(0.6), Inches(1.5), Inches(4), Inches(1.6), fill_color=WHITE, line_color=SURF)
add_text(sl, Inches(0.6), Inches(1.6), Inches(4), Inches(0.25),
         'INVERSION TOTAL', size=10, bold=True, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0.6), Inches(1.9), Inches(4), Inches(0.7),
         '$1.964 MM', size=38, bold=True, color=FOREST, align=PP_ALIGN.CENTER, font_name='Consolas')
add_text(sl, Inches(0.6), Inches(2.6), Inches(4), Inches(0.3),
         '~USD 1,4 M  (TC: 1.400 ARS/USD)', size=11, color=LGRAY, align=PP_ALIGN.CENTER)

# Chart — Inversion donut
fig, ax = plt.subplots(figsize=(4.5, 3.5))
isizes = [1400, 186.2, 345.5, 32.6]
ilabels = ['Nave industrial\n(71%)', 'Obras compl.\n(10%)', 'Maquinaria\n(18%)', 'Vehiculos+\nmob. (2%)']
icolors = [C_FOREST, C_WHEAT, C_NAVY, C_TERRA]
wedges, texts, autotexts = ax.pie(isizes, labels=ilabels, colors=icolors, autopct='%1.0f%%',
    startangle=90, pctdistance=0.75, wedgeprops=dict(width=0.4, edgecolor='white', linewidth=2))
for t in texts: t.set_fontsize(9)
for t in autotexts: t.set_fontsize(9); t.set_color('white'); t.set_weight('bold')
ax.set_title('Distribucion de la inversion', fontsize=11, fontweight='bold', color='#4a4540', pad=12)
chart_buf = make_chart_img(fig)
sl.shapes.add_picture(chart_buf, Inches(0.4), Inches(3.4), Inches(4.5), Inches(3.5))

# Tabla inversion
inv_data = [
    ['Rubro', 'ARS (MM)', '%', 'Detalle'],
    ['Nave industrial', '$1.400,0', '71,3%', 'Planta aislada 1.330 m² + terreno'],
    ['Obras complementarias', '$186,2', '9,5%', 'Electricidad, sanitarias, gas, clima'],
    ['Maquinaria operativa', '$345,5', '17,6%', 'Linea A + Linea B completas'],
    ['Vehiculos y aux.', '$22,8', '1,2%', 'Autoelevador, zorras, balanzas'],
    ['Muebles y utiles', '$9,8', '0,5%', 'PCs (x8), mobiliario, sanitarios'],
    ['TOTAL', '$1.964,3', '100%', ''],
]
add_table(sl, Inches(5.2), Inches(1.5), Inches(7.5), inv_data,
          [Inches(1.8), Inches(1.0), Inches(0.7), Inches(3.0)])

# Maquinaria detalle
add_shape(sl, Inches(5.2), Inches(4.7), Inches(3.6), Inches(2.5), fill_color=WHITE, line_color=SURF)
add_text(sl, Inches(5.4), Inches(4.8), Inches(3.2), Inches(0.3), 'LINEA A — USD equiv.', size=10, bold=True, color=NAVY)
items_a = ['Amasadora: USD 13.500', 'Rotativa: USD 15.000', 'Moldes (x4): USD 20.000',
           'Horno tunel: USD 100.000', 'Enfriamiento: USD 12.000', 'Envasadora: USD 4.000']
add_bullet_frame(sl, Inches(5.4), Inches(5.1), Inches(3.2), Inches(2), items_a, size=9, color=GRAY, spacing=Pt(2))

add_shape(sl, Inches(9.1), Inches(4.7), Inches(3.6), Inches(2.5), fill_color=WHITE, line_color=SURF)
add_text(sl, Inches(9.3), Inches(4.8), Inches(3.2), Inches(0.3), 'LINEA B — ARS/USD', size=10, bold=True, color=TERRA)
items_b = ['Batidora: ARS 2,5M', 'Dosificadora: ARS 2,0M', 'Hornos rot. (x2): ARS 30M',
           'Fermentadora: ARS 1,9M', 'Det. metales: USD 3.200', 'Flow Pack: USD 50.000']
add_bullet_frame(sl, Inches(9.3), Inches(5.1), Inches(3.2), Inches(2), items_b, size=9, color=GRAY, spacing=Pt(2))


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — FLUJO DE FONDOS E INDICADORES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG)
section_header(sl, 9, 'Flujo de Fondos e Indicadores',
    'VAN positivo, TIR >2x costo capital, payback en 3,5 años. El proyecto es viable.')

kpis = [('$1.752M', 'VAN (15%)', FOREST), ('38,85%', 'TIR', NAVY),
        ('3,5 años', 'Payback', TERRA), ('$5.965M', 'Retorno 10a', PLUM)]
for i, (v, l, c) in enumerate(kpis):
    add_kpi(sl, Inches(0.6 + i * 3.05), Inches(1.5), Inches(2.8), v, l, val_color=c)

# Chart — Flujo acumulado
fig, ax = plt.subplots(figsize=(7.5, 3.5))
years = list(range(0, 11))
flujo_neto = [-1964, 793, 793, 793, 793, 793, 793, 793, 793, 793, 793]
acumulado = [-1964, -1171, -378, 414, 1207, 2000, 2793, 3586, 4379, 5172, 5965]
bar_colors = [C_TERRA if f < 0 else C_FOREST for f in flujo_neto]
ax.bar(years, flujo_neto, color=bar_colors, width=0.5, zorder=3, label='Flujo neto anual', alpha=0.8)
ax.plot(years, acumulado, color=C_NAVY, linewidth=2.5, zorder=4, label='Acumulado')
acum_colors = [C_TERRA if a < 0 else C_FOREST for a in acumulado]
for yr, ac, mc in zip(years, acumulado, acum_colors):
    ax.plot(yr, ac, 'o', markersize=7, color=mc, markeredgecolor='white', markeredgewidth=1.5, zorder=5)
ax.axhline(y=0, color='#ccc', linewidth=1, zorder=1)
ax.set_xlabel('Año', fontsize=10)
ax.set_title('Flujo de fondos — horizonte 10 años (ARS MM)', fontsize=11, fontweight='bold', color='#4a4540')
ax.legend(fontsize=9, loc='upper left')
ax.grid(axis='y', alpha=0.2)
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'${int(x):,}'.replace(',','.')))
chart_buf = make_chart_img(fig)
sl.shapes.add_picture(chart_buf, Inches(0.4), Inches(2.9), Inches(8.2), Inches(3.8))

# Criterios
add_shape(sl, Inches(8.8), Inches(2.9), Inches(4.2), Inches(3.8), fill_color=WHITE, line_color=SURF)
add_text(sl, Inches(9.0), Inches(3.0), Inches(3.8), Inches(0.3),
         'CRITERIOS DE DECISION', size=10, bold=True, color=FOREST)
criteria = [
    'VAN > 0→$1.752 MM. El proyecto genera valor por encima del costo de oportunidad.',
    'TIR > WACC→38,85% vs. 15%. Rentabilidad mas del doble de la tasa exigida.',
    'Payback < horizonte→3,5 de 10 años. 6,5 años de flujo positivo neto.',
    'ROI 10 años→~304%. Cada peso invertido genera $3.',
]
add_bullet_frame(sl, Inches(9.0), Inches(3.4), Inches(3.8), Inches(2.8), criteria, size=11, color=GRAY)

# Conclusion box
add_shape(sl, Inches(8.8), Inches(6.0), Inches(4.2), Inches(0.7), fill_color=FOREST)
add_text(sl, Inches(8.8), Inches(6.1), Inches(4.2), Inches(0.5),
         'Proyecto viable bajo los 3\ncriterios estandar de evaluacion.',
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSIONES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG)
section_header(sl, 10, 'Conclusiones',
    'El proyecto es viable desde lo tecnico, comercial y financiero.')

# 3 columnas
cols = [
    ('Oportunidad de mercado', WHEAT, [
        '+1,3M consumidores con demanda insatisfecha.',
        'Competidores artesanales, sin escala.',
        'Precios elevados por falta de oferta industrial.',
        'Tendencia global de crecimiento sostenido.',
    ]),
    ('Ventaja competitiva', FOREST, [
        'Unica planta 100% aislada del pais.',
        'FSSC 22000 + HACCP + ANMAT + Kosher.',
        'Escala industrial = precios competitivos.',
        '7 SKUs veganos, sin alergenos.',
    ]),
    ('Viabilidad financiera', NAVY, [
        'TIR 38,85% (>2x costo capital).',
        'VAN positivo: $1.752 MM.',
        'Payback: 3,5 años.',
        'Retorno 3x en 10 años.',
    ]),
]
for i, (title, color, items) in enumerate(cols):
    x = Inches(0.6 + i * 4.1)
    # Color bar top
    add_shape(sl, x, Inches(1.5), Inches(3.8), Inches(0.08), fill_color=color)
    add_shape(sl, x, Inches(1.58), Inches(3.8), Inches(2.5), fill_color=WHITE, line_color=SURF)
    add_text(sl, x + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.3),
             title, size=12, bold=True, color=color)
    add_bullet_frame(sl, x + Inches(0.2), Inches(2.1), Inches(3.4), Inches(1.8),
                    items, size=11, color=GRAY, spacing=Pt(4))

# Riesgos
risk_data = [
    ['Riesgo critico', 'Mitigacion'],
    ['Demoras SENASA/ANMAT', 'Tramites en paralelo con obra; asesoria regulatoria'],
    ['Cash flow negativo Año 1', 'Capital de trabajo previsto; ramp-up gradual'],
    ['Demora equipos importados', 'Compra anticipada; proveedores alternativos'],
    ['Alta inflacion', 'Clausulas de ajuste; stock estrategico MP'],
]
add_table(sl, Inches(0.6), Inches(4.4), Inches(6.5), risk_data,
          [Inches(2.5), Inches(4.0)], header_color=TERRA)

# Timeline
add_shape(sl, Inches(7.5), Inches(4.4), Inches(5.2), Inches(2.7), fill_color=WHITE, line_color=SURF)
add_text(sl, Inches(7.7), Inches(4.5), Inches(4.8), Inches(0.3),
         'TIMELINE DE IMPLEMENTACION', size=10, bold=True, color=FOREST)
tl_items = [
    'Meses 1-4→Planificacion y constitucion legal.',
    'Meses 3-7→Ingenieria, permisos y certificacion ANMAT.',
    'Meses 4-10→Obra civil y construccion de planta.',
    'Meses 6-11→Compra de equipos y validacion.',
    'Meses 9-12→Reclutamiento y capacitacion (50 pos.).',
    'Mes 13+→Pruebas piloto y puesta en marcha.',
]
add_bullet_frame(sl, Inches(7.7), Inches(4.9), Inches(4.8), Inches(2.1),
                tl_items, size=10, color=GRAY, spacing=Pt(2))


# ══════════════════════════════════════════════════════════════
# SLIDE FINAL — CIERRE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, FOREST)

add_text(sl, Inches(0), Inches(2.0), Inches(13.333), Inches(0.5),
         'CONCLUSION DEL ESTUDIO', size=13, bold=True, color=WHEAT, align=PP_ALIGN.CENTER, font_name='Consolas')

add_text(sl, Inches(1), Inches(2.7), Inches(11.333), Inches(1),
         'El proyecto es viable desde lo tecnico,\ncomercial y financiero.',
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, Inches(1), Inches(4.2), Inches(11.333), Inches(0.5),
         'TIR 38,85%   ·   VAN $1.752 MM   ·   Payback 3,5 años   ·   Retorno 3x',
         size=18, color=RGBColor(0xcc, 0xcc, 0xcc), align=PP_ALIGN.CENTER, font_name='Consolas')

add_text(sl, Inches(0), Inches(6.3), Inches(13.333), Inches(0.3),
         'Proyecto Final — UTN FRBA 2026 — Grupo 8', size=12, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(6.65), Inches(13.333), Inches(0.3),
         'Albarracin  ·  Carrizo  ·  Greco  ·  Rodriguez  ·  Schaltegger', size=11, color=LGRAY, align=PP_ALIGN.CENTER)


# ─── GUARDAR ───
output_path = os.path.join(OUT_DIR, 'Planta_SinTACC_Presentacion_G8.pptx')
prs.save(output_path)
print(f'Presentacion guardada en: {output_path}')
